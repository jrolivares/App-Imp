"""
WhatsApp chat parser + PPTX generator for Mondelez Milka implementations.
"""
import re, json, os, copy, zipfile, csv, difflib, unicodedata, urllib.request, io
from datetime import datetime
from pathlib import Path
from collections import defaultdict
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from lxml import etree


# ── Store database (Google Sheets) ───────────────────────────────────────────

# URL del CSV publicado – se puede sobreescribir con la variable de entorno STORE_DB_URL
SHEETS_CSV_URL = os.environ.get(
    'STORE_DB_URL',
    'https://docs.google.com/spreadsheets/d/e/'
    '2PACX-1vRX_MbNxlPJqTpAg89E51WOrp-oqNd6fAwjlN00ON6-SG1tGzQZNj7ZTs-0vgRAy53u0Dqjhi0I6Cyn'
    '/pub?output=csv'
)

_store_db_cache: list = None      # filas crudas
_store_db_index: list = None      # (nombre_norm, words_set, row) pre-calculado
_lookup_cache:   dict = {}        # query → resultado cacheado


def _norm(text: str) -> str:
    """Normaliza texto: minúsculas, sin tildes, sin puntuación."""
    text = text.lower().strip()
    text = unicodedata.normalize('NFD', text)
    text = ''.join(c for c in text if unicodedata.category(c) != 'Mn')
    text = re.sub(r'[^\w\s]', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def load_store_db() -> list:
    """Carga y cachea la base de tiendas desde Google Sheets CSV."""
    global _store_db_cache, _store_db_index
    if _store_db_cache is not None:
        return _store_db_cache
    try:
        req = urllib.request.urlopen(SHEETS_CSV_URL, timeout=15)
        content = req.read().decode('utf-8-sig')
        rows = list(csv.reader(content.splitlines()))
        _store_db_cache = rows[1:]
        # Pre-normalizar una sola vez al cargar
        # Índice: (nombre_norm, words_set, codigo_upper, row)
        _store_db_index = [
            (_norm(row[5]), set(_norm(row[5]).split()), row[2].strip().upper() if len(row) > 2 else '', row)
            for row in _store_db_cache if len(row) >= 9
        ]
        print(f'[store_db] {len(_store_db_cache)} tiendas cargadas.', flush=True)
    except Exception as exc:
        print(f'[store_db] No se pudo cargar la base de tiendas: {exc}', flush=True)
        _store_db_cache = []
        _store_db_index = []
    return _store_db_cache


def _make_result(row, score):
    return {
        'cadena':      row[3],
        'nombre_sala': row[5],
        'comuna':      row[7],
        'region':      row[8],
        'score':       round(score, 3),
    }


def lookup_store(query: str, chain: str = None, code: str = None):
    """
    Busca la tienda en el índice pre-normalizado.
    Estrategia (en orden de prioridad):
      1. Match exacto por código (col C) — máxima precisión para "local NNN"
      2. Fuzzy match por nombre (col F) con filtro de cadena
    - chain: cadena WhatsApp para filtrar filas por prefijo de nombre
    - code:  código parseado (ej. 'H620', 'N123') para match directo
    """
    cache_key = f"{chain}|{code}|{query}"
    if cache_key in _lookup_cache:
        return _lookup_cache[cache_key]

    load_store_db()
    if not _store_db_index:
        _lookup_cache[cache_key] = None
        return None

    chain_prefix = _norm(chain)[:4] if chain else None

    # ── 1. Buscar por código exacto ───────────────────────────────────────────
    if code:
        code_up   = code.upper().strip()                      # "H066"
        code_bare = re.sub(r'^[A-Z]+', '', code_up)           # "066"
        code_nz   = code_bare.lstrip('0') or '0'              # "66"  (sin ceros iniciales)
        code_up_nz = re.sub(r'^([A-Z]+)', r'\1', code_up).rstrip('0').rstrip() # fallback
        # Set de todas las variantes a probar
        code_variants = {code_up, code_bare, code_nz,
                         re.sub(r'(?<=[A-Z])0+', '', code_up)}  # "H066"→"H66"
        for nombre, nombre_words, db_code, row in _store_db_index:
            if chain_prefix and not nombre.startswith(chain_prefix):
                continue
            if db_code and db_code in code_variants:
                result = _make_result(row, 1.0)
                _lookup_cache[cache_key] = result
                print(f'[lookup] código exacto {code_up} → {row[5]}', flush=True)
                return result

    # ── 2. Fuzzy match por nombre ─────────────────────────────────────────────
    q = _norm(query)
    q_words = set(q.split())
    best_score, best_row = 0.0, None

    for nombre, nombre_words, db_code, row in _store_db_index:
        if chain_prefix and not nombre.startswith(chain_prefix):
            continue
        overlap = len(q_words & nombre_words) / max(len(q_words), 1)
        if overlap > 0.3:
            ratio = difflib.SequenceMatcher(None, q, nombre).ratio()
        else:
            ratio = 0.0
        score = max(ratio, overlap * 0.88)
        if score > best_score:
            best_score = score
            best_row = row

    THRESHOLD = 0.55
    result = None
    if best_score >= THRESHOLD and best_row is not None:
        result = _make_result(best_row, best_score)
    _lookup_cache[cache_key] = result
    return result

CHAIN_ORDER = [
    'SISA', 'JUMBO', 'HIPER', 'SANTA ISABEL', 'TOTTUS', 'SMU', 'UNIMARC',
    'EASY', 'SODIMAC', 'LIDER', 'WALMART', 'ACUENTA',
    'ALVI', 'CASANOVA', 'CENTRAL MAYORISTA', 'COMERCIAL CASTRO', 'CONSTRUMART',
    'CORONA', 'CRUZ VERDE', 'CUGAT', 'DIMARC', 'EKONO', 'ELTIT',
    'FALABELLA', 'FASA', 'HITES', 'KUNCAR',
    'LA MUNDIAL', 'LA OFERTA', 'LA POLAR', 'LIQUIMAX', 'M10', 'MAICAO',
    'OK MARKET', 'OXXO', 'PARIS', 'PREUNIC', 'PROVIMARKET',
    'RIPLEY', 'SALCOBRAND', 'SUDAMERICANA',
    'SUPER 10', 'SUPER OFERTA', 'TALEB', 'TEBA EXPRESS',
]

CHAIN_RULES = [
    # ── Cadenas originales ────────────────────────────────────────────────────
    ('SISA',             r'\bSisa\b|\bSISA\b',                                   'N'),
    ('JUMBO',            r'\bJUMBO\b|\bJumbo\b',                                 'J'),
    ('HIPER',            r'\bHIPER\b|\bHiper\b',                                 'H'),
    ('SANTA ISABEL',     r'Santa\s+Isabel|SANTA\s+ISABEL',                      'SI'),
    ('TOTTUS',           r'\bTOTTUS\b|\bTottus\b',                              'T'),
    ('SMU',              r'\bSMU\b',                                              'S'),
    ('UNIMARC',          r'\bUnimarc\b|\bUNIMARC\b',                            'U'),
    ('EASY',             r'\bEasy\b|\bEASY\b',                                   'E'),
    ('SODIMAC',          r'\bSodimac\b|\bSODIMAC\b|\bHomecenter\b|\bHOMECENTER\b','SO'),
    ('LIDER',            r'\bL[ií]der\b|\bLIDER\b',                              'L'),
    ('WALMART',          r'\bWalmart\b|\bWALMART\b',                              'W'),
    ('ACUENTA',          r'\bAcuenta\b|\bACUENTA\b',                              'AC'),
    # ── Nuevas cadenas ────────────────────────────────────────────────────────
    ('ALVI',             r'\bAlvi\b|\bALVI\b',                                   'AL'),
    ('CASANOVA',         r'\bCasanova\b|\bCASANOVA\b',                          'CAS'),
    ('CENTRAL MAYORISTA',r'Central\s+Mayorista|CENTRAL\s+MAYORISTA',            'CM'),
    ('COMERCIAL CASTRO', r'Comercial\s+Castro|COMERCIAL\s+CASTRO',              'CC'),
    ('CONSTRUMART',      r'\bConstrumart\b|\bCONSTRUMART\b',                    'CON'),
    ('CORONA',           r'\bCorona\b|\bCORONA\b',                              'COR'),
    ('CRUZ VERDE',       r'Cruz\s+Verde|CRUZ\s+VERDE',                          'CV'),
    ('CUGAT',            r'\bCugat\b|\bCUGAT\b',                                'CUG'),
    ('DIMARC',           r'\bDimarc\b|\bDIMARC\b',                              'DIM'),
    ('EKONO',            r'\bEkono\b|\bEKONO\b',                                'EK'),
    ('ELTIT',            r'\bEltit\b|\bELTIT\b',                                'ELT'),
    ('FALABELLA',        r'\bFalabella\b|\bFALABELLA\b',                        'FAL'),
    ('FASA',             r'\bFasa\b|\bFASA\b',                                   'FAS'),
    ('HITES',            r'\bHites\b|\bHITES\b',                                'HIT'),
    ('KUNCAR',           r'\bKuncar\b|\bKUNCAR\b',                              'KUN'),
    ('LA MUNDIAL',       r'La\s+Mundial|LA\s+MUNDIAL',                          'LM'),
    ('LA OFERTA',        r'La\s+Oferta|LA\s+OFERTA',                            'LO'),
    ('LA POLAR',         r'La\s+Polar|LA\s+POLAR',                              'LP'),
    ('LIQUIMAX',         r'\bLiquimax\b|\bLIQUIMAX\b',                          'LIQ'),
    ('M10',              r'\bM10\b',                                              'M10'),
    ('MAICAO',           r'\bMaicao\b|\bMAICAO\b',                              'MAI'),
    ('OK MARKET',        r'Ok\s+Market|OK\s+MARKET|\bOkmarket\b',               'OK'),
    ('OXXO',             r'\bOxxo\b|\bOXXO\b',                                   'OXX'),
    ('PARIS',            r'\bParis\b|\bPARIS\b',                                'PAR'),
    ('PREUNIC',          r'\bPreunic\b|\bPREUNIC\b',                            'PRE'),
    ('PROVIMARKET',      r'\bProvimarket\b|\bPROVIMARKET\b',                    'PRO'),
    ('RIPLEY',           r'\bRipley\b|\bRIPLEY\b',                              'RIP'),
    ('SALCOBRAND',       r'\bSalcobrand\b|\bSALCOBRAND\b',                      'SAL'),
    ('SUDAMERICANA',     r'\bSudamericana\b|\bSUDAMERICANA\b',                  'SUD'),
    ('SUPER 10',         r'Super\s*10|SUPER\s*10',                              'S10'),
    ('SUPER OFERTA',     r'Super\s+Oferta|SUPER\s+OFERTA',                      'SOF'),
    ('TALEB',            r'\bTaleb\b|\bTALEB\b',                                'TAL'),
    ('TEBA EXPRESS',     r'Teba\s+Express|TEBA\s+EXPRESS|\bTeba\b|\bTEBA\b',    'TE'),
]

BAD_WORDS = [
    # Operativos / logística
    'implementación', 'implementacion', 'botadero', 'payloader', 'payloder',
    'material', 'ingreso', 'autorizar', 'autorización', 'autorizacion',
    'abastece', 'abastecer', 'stock',
    'bodega', 'armar', 'solicita', 'solicitud', 'encargada', 'reponedor',
    'rechaza', 'pendiente', 'dejar', 'dejar en', 'nota en',
    'espera', 'esperando',
    # Docs / instructivos
    'instructivo', 'manual', 'capacitación', 'capacitacion', 'pdf',
    'páginas', 'paginas', 'documento',
    # Comunicación / saludo
    'campaña', 'correo', 'problema', 'aparece', 'quieren', 'fecha de término',
    'agregué', 'buen día', 'estará', 'porfa', 'enviados', 'si no',
    'revisando', 'añadir', 'están', 'terminó', 'buen',
    # Señales de que es un reporte, no un identificador de tienda
    'imagen omitida', 'adjunto:', '@',
]

# Soporta formato 12h con AM/PM y formato 24h sin AM/PM
LINE_RE = re.compile(r'^\[(\d{2}-\d{2}-\d{2}), (\d{1,2}:\d{2}:\d{2})(?:\u202f([AP]M))?\] ([^:]+): (.*)$')
# Adjuntos formato antiguo: <attached: file.jpg> / <adjunto: file.jpg>
ATTACH_RE = re.compile(
    r'<(?:attached|adjunto):\s*([^\s>][^>]*\.(?:jpg|jpeg|png|webp))\s*>',
    re.IGNORECASE
)
# Adjuntos formato nuevo: file.jpg (file attached) / file.jpg (archivo adjunto)
ATTACH_RE_PAREN = re.compile(
    r'([\w][\w\-\.]*\.(?:jpg|jpeg|png|webp))\s*\((?:file attached|archivo adjunto)\)',
    re.IGNORECASE
)


def find_photos(text: str) -> list:
    """Extrae nombres de archivo de fotos de un fragmento de texto WhatsApp.
    Soporta formato antiguo (<adjunto: ...>) y formato nuevo (... (archivo adjunto)).
    """
    return ATTACH_RE.findall(text) + ATTACH_RE_PAREN.findall(text)


# ── Chat parsing ──────────────────────────────────────────────────────────────

def parse_messages(chat_text: str) -> list:
    content = chat_text.replace('\u200e', '').replace('\u200f', '').replace('\r', '')
    messages, current = [], None
    for line in content.split('\n'):
        m = LINE_RE.match(line)
        if m:
            if current:
                messages.append(current)
            d, t, ap, sender, txt = m.groups()
            if ap:  # formato 12h con AM/PM
                dt = datetime.strptime(f'{d} {t} {ap}', '%d-%m-%y %I:%M:%S %p')
            else:   # formato 24h sin AM/PM
                dt = datetime.strptime(f'{d} {t}', '%d-%m-%y %H:%M:%S')
            current = {'dt': dt, 'sender': sender.strip(), 'text': txt,
                       'photos': find_photos(txt)}
        else:
            if current:
                current['text'] += '\n' + line
                current['photos'] += find_photos(line)
    if current:
        messages.append(current)
    return messages


def detect_chain(text):
    first = text.strip().split('\n')[0]
    for chain, pat, prefix in CHAIN_RULES:
        if re.search(pat, first, re.IGNORECASE):
            return chain, prefix
    return None, None


def is_store_message(text):
    chain, _ = detect_chain(text)
    if not chain:
        return False
    lines = text.strip().split('\n')
    first = lines[0].lower()

    # Siempre rechazar si la PRIMERA línea tiene palabras de reporte/logística
    if any(b in first for b in BAD_WORDS):
        return False

    # Calcular qué queda de la primera línea después de quitar el nombre de cadena
    chain_stripped = re.sub(
        r'\b(?:sisa|jumbo|hiper|lider|líder|santa\s+isabel|tottus|smu|unimarc|easy|'
        r'sodimac|homecenter|walmart|acuenta|alvi|casanova|central\s+mayorista|'
        r'comercial\s+castro|construmart|corona|cruz\s+verde|cugat|dimarc|ekono|eltit|'
        r'falabella|fasa|hites|kuncar|la\s+mundial|la\s+oferta|la\s+polar|liquimax|'
        r'm10|maicao|ok\s+market|oxxo|paris|preunic|provimarket|ripley|salcobrand|'
        r'sudamericana|super\s*10|super\s+oferta|taleb|teba(?:\s+express)?)\b',
        '', first, flags=re.IGNORECASE
    ).strip()

    # Solo revisar la SEGUNDA línea para bad words si la primera línea era únicamente
    # el nombre de la cadena (sin dirección ni código), como "JUMBO\nEn capacitación..."
    if len(lines) > 1 and not chain_stripped:
        if any(b in lines[1].lower() for b in BAD_WORDS):
            return False

    # Rechazar patrón "CADENA - [texto libre]" → es un reporte, no identificador
    if chain_stripped.startswith('- ') or chain_stripped.startswith('-\t'):
        return False

    return True


def parse_store_line(text, chain, prefix):
    first = re.sub(r'<(?:attached|adjunto):[^>]+>', '', text.strip().split('\n')[0]).strip()
    parts = first.split('\t')
    if len(parts) >= 3:
        raw_code = re.sub(r'[°º]', '', parts[0].strip())
        addr = parts[2].strip()
        city = parts[3].strip() if len(parts) > 3 else ''
        if re.match(r'^[A-Za-z]\d+$', raw_code):
            code = raw_code.upper()
        elif re.match(r'^\d+$', raw_code):
            code = prefix + raw_code
        else:
            code = raw_code.upper()
        return code, addr, city
    # Free-form
    stripped = first
    for _, pat, _ in CHAIN_RULES:
        stripped = re.sub(pat, '', stripped, flags=re.IGNORECASE).strip(' ,\t')
    # Patrón numérico al inicio: "123 Dirección"
    m = re.match(r'^(\d+)\s+(.+)$', stripped)
    if m:
        return prefix + m.group(1), m.group(2).strip(), ''
    # Patrón "local NNN" o "líder local NNN" (usado por HIPER y otros)
    m_local = re.search(r'\blocal[:\s]*(\d+)\b', stripped, re.IGNORECASE)
    if m_local:
        return prefix + m_local.group(1), stripped.strip(), ''
    return None, stripped.strip(), ''


def parse_status(text):
    pl, bt, notes = None, 0, []
    for line in text.split('\n'):
        l, lu = line.strip(), line.strip().upper()
        if '\u2705' in l:
            if 'PAYLOAD' in lu:
                pl = 'Implementado'
            if 'BOTADERO' in lu:
                mm = re.search(r'(\d+)\s*BOTADERO', lu)
                bt = int(mm.group(1)) if mm else 1
        if '\u274c' in l:
            if 'PAYLOAD' in lu:
                pl = 'No implementado'
            note = re.sub(r'[\u274c\u2705]', '', l).strip()
            if note:
                notes.append(note)
        if 'NO SE PUDO IMPLEMENTAR' in lu and 'PAYLOAD' in lu:
            pl = 'No implementado'
    return pl, bt, ' | '.join(notes)


def extract_stores(messages: list, start_date: datetime, end_date: datetime) -> list:
    recent = [m for m in messages if start_date <= m['dt'] < end_date]

    # ── Fase 1: identificar todos los mensajes de tienda ──────────────────────
    store_entries = []   # (msg_index, chain, prefix, code, address, city, msg)
    for i, msg in enumerate(recent):
        if is_store_message(msg['text']):
            chain, prefix = detect_chain(msg['text'])
            code, address, city = parse_store_line(msg['text'], chain, prefix)
            store_entries.append((i, chain, prefix, code, address, city, msg))

    # ── Fase 2: asignar fotos a tiendas ───────────────────────────────────────
    store_photos = defaultdict(list)   # store_entry_index → [filenames]
    photo_timestamps: dict = {}        # filename → datetime del mensaje que la contenía

    # Primero: fotos dentro del MISMO mensaje de tienda (enviadas juntas)
    # Bug histórico: el loop de abajo hace break cuando i >= j, así que
    # si la foto y la tienda están en el mismo mensaje (i==j) se perdía.
    store_msg_indices = {i: k for k, (i, *_) in enumerate(store_entries)}
    for j, msg in enumerate(recent):
        if msg['photos'] and j in store_msg_indices:
            for p in msg['photos']:
                photo_timestamps[p] = msg['dt']
            store_photos[store_msg_indices[j]].extend(msg['photos'])

    # Luego: fotos enviadas ANTES o DESPUÉS del mensaje de tienda
    unassigned = 0
    for j, msg in enumerate(recent):
        if not msg['photos']:
            continue
        if j in store_msg_indices:
            continue   # ya capturadas arriba

        msg_dt = msg['dt']
        best_same_sender    = None   # foto DESPUÉS del ID, mismo sender  (8h)
        best_any_sender     = None   # foto DESPUÉS del ID, cualquier sender (1h)
        best_forward_sender = None   # foto ANTES del ID, mismo sender (5min)
        # ↑ Cubre el patrón "Romina manda foto y luego escribe el nombre de la tienda"

        for k, (i, chain, prefix, code, address, city, smsg) in enumerate(store_entries):
            if i < j:
                # Tienda ANTES de la foto (lógica original)
                diff = (msg_dt - smsg['dt']).total_seconds()
                if diff < 0:
                    continue
                if smsg['sender'] == msg['sender']:
                    if diff <= 28800:
                        if best_same_sender is None or diff < best_same_sender[0]:
                            best_same_sender = (diff, k)
                if diff <= 3600:
                    if best_any_sender is None or diff < best_any_sender[0]:
                        best_any_sender = (diff, k)
            else:
                # Tienda DESPUÉS de la foto (ventana hacia adelante)
                diff_ahead = (smsg['dt'] - msg_dt).total_seconds()
                if diff_ahead > 300:   # más de 5 min → cortar búsqueda
                    break
                if smsg['sender'] == msg['sender'] and diff_ahead >= 0:
                    if best_forward_sender is None or diff_ahead < best_forward_sender[0]:
                        best_forward_sender = (diff_ahead, k)

        # Prioridades: mismo sender (hacia atrás) > mismo sender (hacia adelante) > cualquier sender
        target = best_same_sender or best_forward_sender or best_any_sender
        if target:
            for p in msg['photos']:
                photo_timestamps[p] = msg['dt']
            store_photos[target[1]].extend(msg['photos'])
        else:
            unassigned += 1

    print(f'[extract] fotos sin tienda asignada: {unassigned}', flush=True)

    # ── Fase 3: construir registros de tienda ─────────────────────────────────
    raw = []
    for k, (i, chain, prefix, code, address, city, msg) in enumerate(store_entries):
        photos = store_photos[k]
        seen = set()
        photos = [p for p in photos if not (p in seen or seen.add(p))]
        pl, bt, notes = parse_status(msg['text'])
        print(f'[extract] tienda={code or address[:30]!r} fotos={len(photos)}', flush=True)

        # Buscar tienda en base de datos formal (código primero, luego fuzzy por nombre)
        query = f"{chain} {address}" if address else chain
        db_match = lookup_store(query, chain=chain, code=code)

        raw.append({
            'chain': chain, 'code': code, 'address': address, 'city': city,
            'sender': msg['sender'], 'date': msg['dt'].strftime('%d/%m/%Y'),
            'datetime': msg['dt'].isoformat(),
            'photos': photos, 'payloader': pl, 'botaderos': bt, 'notes': notes,
            # Timestamp por foto: filename → datetime (para mostrar la fecha real en el caption)
            'photo_timestamps': {p: photo_timestamps[p] for p in photos if p in photo_timestamps},
            # Datos formales desde la planilla (None si no hubo match)
            'db_cadena':      db_match['cadena']      if db_match else None,
            'db_nombre_sala': db_match['nombre_sala']  if db_match else None,
            'db_comuna':      db_match['comuna']       if db_match else None,
            'db_region':      db_match['region']       if db_match else None,
        })
    # Deduplicate: misma cadena + mismo nombre (oficial si hay, si no parseado) + misma fecha
    deduped = {}
    for s in raw:
        name_key = s.get('db_nombre_sala') or s['code'] or s['address'][:25]
        key = f"{s['chain']}_{name_key}_{s['date']}"
        if key in deduped:
            ex = deduped[key]
            merged = ex['photos'] + s['photos']
            seen = set()
            ex['photos'] = [p for p in merged if not (p in seen or seen.add(p))]
            ex.setdefault('photo_timestamps', {}).update(s.get('photo_timestamps', {}))
            if s['payloader']:
                ex['payloader'] = s['payloader']
            if s['botaderos']:
                ex['botaderos'] = s['botaderos']
        else:
            deduped[key] = s
    stores = list(deduped.values())
    stores.sort(key=lambda s: (
        CHAIN_ORDER.index(s['chain']) if s['chain'] in CHAIN_ORDER else 99,
        s['datetime'],
        s['code'] or ''
    ))
    return stores


# ── PPTX generation ───────────────────────────────────────────────────────────

def open_corrected(img_path: str, max_px: int = 1200):
    """
    Corrige orientación EXIF y reduce resolución solo si es necesario.
    - Si la foto ya es pequeña y no está rotada → devuelve None (usa archivo original, rápido).
    - Si necesita ajuste → devuelve BytesIO con imagen corregida.
    """
    try:
        img = Image.open(img_path)   # lazy: no decodifica píxeles todavía

        # Leer orientación EXIF sin decodificar la imagen completa
        try:
            orientation = (img.getexif() or {}).get(274, 1)
        except Exception:
            orientation = 1

        needs_rotate = orientation not in (1, 0, None)
        needs_resize = img.width > max_px or img.height > max_px

        if not needs_rotate and not needs_resize:
            return None   # sin cambios → usar archivo original directamente

        img.load()   # decodificar solo si realmente hay que modificar
        img = ImageOps.exif_transpose(img)
        if img.mode not in ('RGB', 'L'):
            img = img.convert('RGB')
        if img.width > max_px or img.height > max_px:
            img.thumbnail((max_px, max_px), Image.LANCZOS)

        buf = io.BytesIO()
        img.save(buf, format='JPEG', quality=82)
        buf.seek(0)
        return buf
    except Exception:
        return None


def _variance_stats(img):
    """
    Returns (mean_row_var, mean_col_var) on a 96×72 grayscale thumbnail.
    Row variance = horizontal variation within each row.
    Col variance = vertical variation within each column.
    """
    thumb = img.convert('L').resize((96, 72), Image.BOX)
    tw, th = thumb.size
    px = list(thumb.getdata())
    total_row_var = 0.0
    for r in range(th):
        row = px[r * tw:(r + 1) * tw]
        m = sum(row) / tw
        total_row_var += sum((p - m) ** 2 for p in row) / tw
    total_col_var = 0.0
    for c in range(tw):
        col = px[c::tw]
        m = sum(col) / th
        total_col_var += sum((p - m) ** 2 for p in col) / th
    return total_row_var / th, total_col_var / tw


def fit_photo_to_slot(img_path: str, slot_w_emu: int, slot_h_emu: int, max_px: int = 1200) -> io.BytesIO:
    """
    Prepara una foto para insertarse en un slot del PPTX:
    1. Corrige orientación EXIF.
    2. Heurística de contenido: si columnas >> filas en varianza → foto portrait
       almacenada landscape → rota 90° CCW. Funciona independientemente del aspect ratio.
    3. Center-crop al mismo aspect ratio del slot (evita distorsión).
    4. Reduce resolución si es necesario.
    Siempre devuelve un BytesIO listo para usar.
    """
    fname = os.path.basename(img_path)
    try:
        img = Image.open(img_path)
        try:
            exif          = img.getexif() or {}
            exif_has_data = len(exif) > 0
            orientation   = exif.get(274, 1)
        except Exception:
            exif_has_data = False
            orientation   = 1
        img.load()
        img = ImageOps.exif_transpose(img)
        if img.mode not in ('RGB', 'L'):
            img = img.convert('RGB')

        w0, h0 = img.size

        # ── Diagnóstico: siempre loguear para depuración ──────────────────────────
        try:
            rv, cv = _variance_stats(img)
            ratio_vc = cv / rv if rv > 0 else 0
            print(f'[img-diag] {fname} size={w0}x{h0} exif={exif_has_data}'
                  f' orient={orientation} row_var={rv:.1f} col_var={cv:.1f}'
                  f' col/row={ratio_vc:.2f}', flush=True)
        except Exception as diag_err:
            print(f'[img-diag] {fname} diag-error: {diag_err}', flush=True)

        # ── Heurística de rotación basada en contenido ────────────────────────────
        # WhatsApp elimina el tag de orientación (274) sin rotar los píxeles.
        # Detectamos el contenido mal-orientado comparando varianza de filas vs columnas.
        #
        # CASO A – landscape almacenada como portrait (w > h):
        #   Contenido portrait (persona de pie) → estructuras verticales fuertes
        #   → col_var >> row_var → col/row > 2.5 → rotar 90° CCW
        #
        # CASO B – portrait almacenada como landscape rotada (h > w):
        #   Foto tomada landscape (estanterías horizontales) pero almacenada portrait
        #   con los píxeles rotados 90° → estanterías aparecen verticales en la imagen
        #   → row_var >> col_var → col/row < 1.0 → rotar 90° CW
        #   (Problema confirmado con Xiaomi Redmi Note 11 Pro 5G via HyperOS/WhatsApp)
        #
        # Umbral col/row < 1.0 confirmado empíricamente: rotadas = 0.39–0.97, correctas ≥ 1.10
        if orientation == 1:
            try:
                rv, cv = _variance_stats(img)
                ratio_vc = cv / rv if rv > 0 else 0
                if w0 > h0 and ratio_vc > 2.5:
                    # Caso A: portrait content stored landscape → rotate CCW
                    img = img.rotate(90, expand=True)
                    print(f'[img] rotated CCW (case A): {fname}'
                          f' ({w0}x{h0}→{img.size[0]}x{img.size[1]})'
                          f' col/row={ratio_vc:.2f}', flush=True)
                elif h0 > w0 and ratio_vc < 1.0:
                    # Caso B: landscape content stored portrait → rotate CCW
                    img = img.rotate(90, expand=True)
                    print(f'[img] rotated CCW (case B): {fname}'
                          f' ({w0}x{h0}→{img.size[0]}x{img.size[1]})'
                          f' col/row={ratio_vc:.2f}', flush=True)
            except Exception:
                pass

        # Center-crop al ratio del slot
        target_ratio = slot_w_emu / slot_h_emu if slot_h_emu else 1
        iw, ih = img.size
        img_ratio = iw / ih if ih else 1

        if abs(img_ratio - target_ratio) > 0.05:   # diferencia mayor al 5%
            if img_ratio > target_ratio:
                # foto más ancha que el slot → recortar lados
                new_w = int(ih * target_ratio)
                x0 = (iw - new_w) // 2
                img = img.crop((x0, 0, x0 + new_w, ih))
            else:
                # foto más alta que el slot → recortar arriba/abajo
                new_h = int(iw / target_ratio)
                y0 = (ih - new_h) // 2
                img = img.crop((0, y0, iw, y0 + new_h))

        if img.width > max_px or img.height > max_px:
            img.thumbnail((max_px, max_px), Image.LANCZOS)

        buf = io.BytesIO()
        img.save(buf, format='JPEG', quality=82)
        buf.seek(0)
        return buf
    except Exception:
        return None


def build_photo_index(photos_dir: str) -> dict:
    """
    Construye un índice {nombre_archivo → path_completo} buscando recursivamente
    en photos_dir. Necesario cuando el ZIP tiene subcarpetas (ej. Media/).
    """
    index = {}
    for root, _, files in os.walk(photos_dir):
        for f in files:
            if f.lower().endswith(('.jpg', '.jpeg', '.png', '.webp')):
                # Guardamos solo si no está ya (prioridad a rutas más cortas / raíz)
                if f not in index:
                    index[f] = os.path.join(root, f)
    return index


def select_photos(photos, n, photos_dir, _photo_index=None):
    """
    Selecciona hasta n fotos de la lista y devuelve sus paths completos.
    Usa _photo_index si ya fue construido; si no, lo construye internamente.
    """
    if n == 0:
        return []
    idx = _photo_index if _photo_index is not None else build_photo_index(photos_dir)
    avail = [idx[p] for p in photos if p in idx]
    if not avail:
        return []
    if len(avail) <= n:
        return avail
    step = len(avail) / n
    return [avail[int(i * step)] for i in range(n)]


def group_runs_by_br(para):
    groups, run_idx = [[]], 0
    for child in para._p:
        tag = etree.QName(child).localname
        if tag == 'r':
            groups[-1].append(run_idx)
            run_idx += 1
        elif tag == 'br':
            groups.append([])
    return groups


def set_line(para, idxs, text):
    runs = para.runs
    if not idxs:
        return
    if idxs[0] < len(runs):
        runs[idxs[0]].text = text
    for i in idxs[1:]:
        if i < len(runs):
            runs[i].text = ''


def update_caption(shape, photo_date, slot=0, pl_status=None, bt_status=None):
    """Actualiza el caption de una foto con la fecha real de envío.

    Detecta automáticamente el formato del template:
    - Simple (1 línea "FECHA: ..."): solo actualiza la fecha → para el nuevo Termplate
    - Multi-línea (FOTO/FECHA/ELEMENTO/STATUS): rellena todas las líneas → template legado
    """
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    groups = group_runs_by_br(para)
    n = len(groups)

    if n <= 1:
        # Formato simple: solo hay una línea (ej. "FECHA: 11/03/2026")
        # Reemplazar con la fecha real de la foto
        set_line(para, groups[0] if groups else [], f'FECHA: {photo_date}')
    else:
        # Formato multi-línea (template legado con FOTO/FECHA/ELEMENTO/STATUS)
        element = 'Payloader Easter' if slot == 0 else 'Botadero Easter'
        status = (pl_status if slot == 0 else bt_status) or ''
        if n >= 1:
            set_line(para, groups[0], f'FOTO IMPLEMENTACIÓN {slot + 1}')
        if n >= 2:
            set_line(para, groups[1], f'FECHA: {photo_date}')
        if n >= 3:
            g, runs = groups[2], para.runs
            if len(g) >= 3:
                runs[g[0]].text = 'ELEMENTO: '
                runs[g[1]].text = 'Payloader' if slot == 0 else 'Botadero'
                runs[g[2]].text = ' Easter'
                for x in g[3:]:
                    runs[x].text = ''
            elif g:
                runs[g[0]].text = f'ELEMENTO: {element}'
                for x in g[1:]:
                    runs[x].text = ''
        if n >= 4:
            set_line(para, groups[3], f'STATUS: {status}')


_NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def _collect_pic_slots(slide):
    """Return absolute (left, top, w, h) for every picture in the slide,
    including pictures inside group shapes.  Sorted reading-order (top→bottom,
    left→right).
    """
    slots = []
    for shape in slide.shapes:
        if shape.shape_type == 13:                          # direct PICTURE
            slots.append((shape.left, shape.top, shape.width, shape.height))
        elif shape.shape_type == 6:                         # GROUP
            grp = shape._element
            xfrm = grp.find(f'.//{{{_NS_A}}}xfrm')
            if xfrm is None:
                continue
            off   = xfrm.find(f'{{{_NS_A}}}off')
            ext   = xfrm.find(f'{{{_NS_A}}}ext')
            chOff = xfrm.find(f'{{{_NS_A}}}chOff')
            chExt = xfrm.find(f'{{{_NS_A}}}chExt')
            if any(e is None for e in [off, ext, chOff, chExt]):
                continue
            gx  = int(off.get('x', 0));  gy  = int(off.get('y', 0))
            gcx = int(ext.get('cx', 1)); gcy = int(ext.get('cy', 1))
            cx0 = int(chOff.get('x', 0)); cy0 = int(chOff.get('y', 0))
            ccx = int(chExt.get('cx', 1)); ccy = int(chExt.get('cy', 1))
            sx = gcx / ccx; sy = gcy / ccy
            for pic in grp.findall(f'{{{_NS_P}}}pic'):
                cxfrm = pic.find(f'.//{{{_NS_A}}}xfrm')
                if cxfrm is None:
                    continue
                co = cxfrm.find(f'{{{_NS_A}}}off')
                ce = cxfrm.find(f'{{{_NS_A}}}ext')
                if co is None or ce is None:
                    continue
                slots.append((
                    int(gx + (int(co.get('x', 0)) - cx0) * sx),
                    int(gy + (int(co.get('y', 0)) - cy0) * sy),
                    int(int(ce.get('cx', 0)) * sx),
                    int(int(ce.get('cy', 0)) * sy),
                ))
    slots.sort(key=lambda s: (round(s[1] / 500000), s[0]))
    return slots


def _remove_all_pics(slide):
    """Remove all PICTURE (13) and GROUP (6) shapes from the slide spTree."""
    spTree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if shape.shape_type in (13, 6):
            spTree.remove(shape._element)


def update_store_slide(slide, store, photos_dir, photo_index=None):
    """Fill a store slide using the new Termplate Sell-Out REV layout."""
    # ── collect slot positions BEFORE removing anything ───────────────────
    pic_slots = _collect_pic_slots(slide)

    code    = store.get('code', '')
    address = store.get('address', '')
    city    = store.get('city', '')
    chain   = store.get('chain', '')
    fecha   = store.get('date', '--/--/----')
    photos  = store.get('photos', [])

    db_nombre = store.get('db_nombre_sala')
    db_comuna = store.get('db_comuna')
    db_region = store.get('db_region')

    if db_nombre:
        header_text = db_nombre
        if db_comuna:
            header_text += f' — {db_comuna}'
        if db_region and db_region != db_comuna:
            header_text += f', {db_region}'
    else:
        chain_label = chain or 'SISA'
        header_text = (f'{code} {chain_label} - {address}' if code
                       else f'{chain_label} - {address}')
        if city:
            header_text += f', {city}'

    text_shapes = [s for s in slide.shapes if s.has_text_frame]

    # ── store-name box: contains 'NOMBRE TIENDA' or is widest non-FECHA box ─
    for sh in sorted(text_shapes, key=lambda s: -s.width):
        t = sh.text_frame.text
        if 'NOMBRE TIENDA' in t or 'FECHA' not in t:
            if sh.text_frame.paragraphs and sh.text_frame.paragraphs[0].runs:
                sh.text_frame.paragraphs[0].runs[0].text = header_text
                for r in sh.text_frame.paragraphs[0].runs[1:]:
                    r.text = ''
            break

    # ── FECHA boxes: one per photo slot, sorted reading-order ────────────────
    fecha_boxes = sorted(
        [s for s in text_shapes if 'FECHA' in s.text_frame.text],
        key=lambda s: (round(s.top / 500000), s.left),
    )

    photo_timestamps = store.get('photo_timestamps', {})
    sel = select_photos(photos, len(pic_slots), photos_dir, _photo_index=photo_index)

    photo_dates = []
    for img_path in sel:
        bn = os.path.basename(img_path)
        dt = photo_timestamps.get(bn)
        photo_dates.append(dt.strftime('%d/%m/%Y') if dt else fecha)

    for i, fbox in enumerate(fecha_boxes):
        if i < len(photo_dates):
            tf = fbox.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = f'FECHA: {photo_dates[i]}'
                for r in tf.paragraphs[0].runs[1:]:
                    r.text = ''
        else:
            try:
                slide.shapes._spTree.remove(fbox._element)
            except Exception:
                pass

    # ── replace picture shapes with real photos ───────────────────────────────
    _remove_all_pics(slide)
    for i, (left, top, w, h) in enumerate(pic_slots):
        if i < len(sel):
            try:
                img_src = fit_photo_to_slot(sel[i], w, h) or open_corrected(sel[i]) or sel[i]
                slide.shapes.add_picture(img_src, left, top, w, h)
            except Exception as e:
                print(f'[pptx] ERROR foto {os.path.basename(sel[i])}: {e}', flush=True)


def add_slide_copy(prs, src_idx):
    src = prs.slides[src_idx]
    new = prs.slides.add_slide(src.slide_layout)
    for shape in list(new.shapes):
        new.shapes._spTree.remove(shape._element)

    # Copiar relaciones de imagen del slide origen al nuevo slide.
    # Sin esto, los r:embed="rIdX" copiados en el XML no tienen contraparte
    # en el .rels del nuevo slide → las imágenes aparecen rotas en PowerPoint.
    rId_map = {}
    for rId, rel in src.part.rels.items():
        if 'image' in rel.reltype:
            new_rId = new.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId

    for shape in src.shapes:
        elem = copy.deepcopy(shape._element)
        if rId_map:
            xml_str = etree.tostring(elem, encoding='unicode')
            for old_rId, new_rId in rId_map.items():
                xml_str = xml_str.replace(f'r:embed="{old_rId}"', f'r:embed="{new_rId}"')
                xml_str = xml_str.replace(f'r:link="{old_rId}"', f'r:link="{new_rId}"')
            elem = etree.fromstring(xml_str)
        new.shapes._spTree.append(elem)
    return new


def make_chain_divider(prs, chain_name, chain_intro_idx=1):
    """Copy the chain-intro template slide (index 1) and fill in the chain name."""
    new = add_slide_copy(prs, chain_intro_idx)
    for sh in new.shapes:
        if sh.has_text_frame and 'NOMBRE CADENA' in sh.text_frame.text:
            tf = sh.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = chain_name
                for r in tf.paragraphs[0].runs[1:]:
                    r.text = ''
            break
    return new


def _detect_template_layout(prs):
    """Auto-detect template structure from slide count.

    Returns (tmpl_count, photo_to_idx, has_closing_slide).

    9-slide template (Template Sell-Out.pptx):
      0=title  1=chain  2=1p  3=2p  4=3p  5=5p  6=4p  7=6p  8=closing

    8-slide template (Termplate Sell-Out REV.pptx):
      0=title  1=chain  2=1p  3=2p  4=3p  5=4p  6=5p  7=6p
    """
    n = len(prs.slides)
    if n >= 9:
        return 9, {1: 2, 2: 3, 3: 4, 4: 6, 5: 5, 6: 7}, True
    return 8, {1: 2, 2: 3, 3: 4, 4: 5, 5: 6, 6: 7}, False


def generate_pptx(stores: list, photos_dir: str, template_path: str, output_path: str) -> dict:
    """Generate the combined PPTX. Returns summary dict.

    Supports both 8-slide and 9-slide templates (auto-detected).
    """
    prs = Presentation(template_path)

    tmpl_count, photo_to_idx, has_closing = _detect_template_layout(prs)

    photo_index = build_photo_index(photos_dir)
    print(f'[pptx] {len(photo_index)} fotos indexadas en {photos_dir}', flush=True)
    print(f'[pptx] template slides={len(prs.slides)} closing={has_closing}', flush=True)

    by_chain = defaultdict(list)
    for s in stores:
        by_chain[s['chain']].append(s)

    # Copy title slide first (template index 0)
    add_slide_copy(prs, 0)

    summary = {}
    for chain in CHAIN_ORDER:
        chain_stores = by_chain.get(chain, [])
        if not chain_stores:
            continue
        make_chain_divider(prs, chain, chain_intro_idx=1)
        for store in chain_stores:
            # Choose template slide by available photo count (1–6)
            avail_count = sum(1 for p in store.get('photos', []) if p in photo_index)
            n = max(1, min(avail_count, 6))
            tmpl_idx = photo_to_idx[n]
            new_slide = add_slide_copy(prs, tmpl_idx)
            update_store_slide(new_slide, store, photos_dir, photo_index=photo_index)
        summary[chain] = len(chain_stores)

    # Closing slide (copy of last template slide, if present)
    if has_closing:
        add_slide_copy(prs, tmpl_count - 1)

    # Remove all original template slides (reverse order to keep indices valid)
    for idx in range(tmpl_count - 1, -1, -1):
        rId = prs.slides._sldIdLst[idx].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[idx]

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return summary


def process_zip(zip_path: str, photos_dir: str, start_date: datetime,
                end_date: datetime, template_path: str, output_path: str) -> dict:
    """Full pipeline: unzip → parse → generate PPTX. Returns result dict."""
    photos_dir = Path(photos_dir)
    photos_dir.mkdir(parents=True, exist_ok=True)

    # Extract ZIP, stripping any top-level subfolder and __MACOSX junk.
    # WhatsApp ZIPs exported from iOS/macOS wrap everything in a subfolder like
    # "WhatsApp Chat - Grupo XYZ/" and add "__MACOSX/" metadata entries.
    with zipfile.ZipFile(zip_path, 'r') as zf:
        for member in zf.infolist():
            # Skip macOS metadata
            if '__MACOSX' in member.filename or member.filename.startswith('.'):
                continue
            parts = Path(member.filename).parts
            # Strip leading subfolder if all files share one (iOS export style)
            stripped = Path(*parts[1:]) if len(parts) > 1 else Path(parts[0])
            target = photos_dir / stripped
            if member.filename.endswith('/'):
                target.mkdir(parents=True, exist_ok=True)
            else:
                target.parent.mkdir(parents=True, exist_ok=True)
                target.write_bytes(zf.read(member.filename))

    # Find _chat.txt (may be at root or one level deep after stripping)
    chat_file = photos_dir / '_chat.txt'
    if not chat_file.exists():
        # Fallback: search recursively
        found = list(photos_dir.rglob('_chat.txt'))
        if found:
            chat_file = found[0]
        else:
            raise FileNotFoundError('No se encontró _chat.txt en el ZIP')

    photos_dir = str(photos_dir)   # rest of pipeline expects str
    chat_text = chat_file.read_text(encoding='utf-8', errors='replace')
    messages = parse_messages(chat_text)
    stores = extract_stores(messages, start_date, end_date)

    if not stores:
        raise ValueError('No se encontraron tiendas en el rango de fechas seleccionado')

    summary = generate_pptx(stores, photos_dir, template_path, output_path)
    total_slides = 1 + sum(len(v) + 1 for v in defaultdict(list,
                           {c: [s for s in stores if s['chain'] == c]
                            for c in CHAIN_ORDER}).values() if v)
    return {
        'stores': stores,
        'summary': summary,
        'total_slides': total_slides,
        'output_path': output_path,
    }
